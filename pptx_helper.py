from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE
from os import path
import os
import shortuuid

#google cloud
from google.cloud import storage
import os
storage_client = storage.Client.from_service_account_json("serviceaccount.json")



#check if user id exists

def create_bucket(userid):
    
    userid = [userid]

    try:

        for name in userid:
  
            BUCKET_NAME = name
            BUCKET = storage_client.bucket(BUCKET_NAME)
            BUCKET = storage_client.create_bucket(BUCKET_NAME, location="us")


            #make objects publicly readable

            """Set a public IAM Policy to bucket"""

            policy = BUCKET.get_iam_policy(requested_policy_version=3)
            policy.bindings.append(
                {"role": "roles/storage.objectViewer", "members": ["allUsers"]}
            )
            BUCKET.set_iam_policy(policy)

            
    except Exception as e:
        print("Unable to create GCP bucket, Error:", e)

    return BUCKET_NAME

def upload_blob(bucket_name, source_file_name, destination_blob_name):

    try:
        storage_client = storage.Client.from_service_account_json("serviceaccount.json")


        bucket = storage_client.bucket(bucket_name)
        blob = bucket.blob(destination_blob_name)

        blob.upload_from_filename(source_file_name)

        print(
            "File {} uploaded to {}.".format(
                source_file_name, destination_blob_name
            )
        )
    
    except Exception as e:
        print("Failed to upload to GCP bucket, Error:", e)


def pptx_extractor(pptx, pptx_filename, TEMP_DIR, BUCKET_NAME, user_id):

    try:
        ppt = Presentation(pptx)

        # upload pptx to GCS bucket

        #uuid in pptx filename
        upload_blob(bucket_name=BUCKET_NAME, 
                source_file_name=pptx, 
                destination_blob_name=f'{user_id}/pptx/pptx/{pptx_filename}_{shortuuid.uuid()}.pptx')

        pptdct = {}

        slidecounter = 1
        for slide in ppt.slides:
            textdct = {}

            textlst = []
            tablelst = []
            imagelst = []
            try: 
                for shape in slide.shapes:

                        # JSON - text
                        if hasattr(shape, "text"):
                            text = shape.text
                            textlst.append(text)
                            text = text.replace('\n',' ')
                            text = text.replace('\x0b',' ')
                            textdct["text"] = textlst

                            json_file_name = (f'{pptx_filename}_{shortuuid.uuid()}.json')

                        if not textlst:
                            textdct["text"] = []

                            json_file_name = (f'{pptx_filename}_{shortuuid.uuid()}.json')

            except Exception as e:
                print("Unable to extract PPT text, Error:", e)
                

            try: 
                for shape in slide.shapes:
                    # Tables - CSV
                    if shape.has_table:
                        table = shape.table
                        df = [['' for i in range(len(table.columns))] for j in range(len(table.rows))]
                        for i, row in enumerate(table.rows):
                            for j, cell in enumerate(row.cells):
                                if cell.text:
                                    df[i][j] = cell.text
                        table_df = pd.DataFrame(df)
                        table_df.rename(columns=table_df.iloc[0], inplace = True)
                        table_df.drop([0], inplace = True)

                        csv_file_name = (f'{pptx_filename}_{shortuuid.uuid()}.csv')

                        # download csv
                        table_df.to_csv(f'{TEMP_DIR}/tables/{csv_file_name}', index=False)

                        # upload csv to GCS bucket
                        upload_blob(bucket_name=BUCKET_NAME, 
                            source_file_name=f"{TEMP_DIR}/tables/{csv_file_name}", 
                            destination_blob_name=f'{user_id}/pptx/tables/{csv_file_name}')

                        tablelst.append(f'https://storage.googleapis.com/{BUCKET_NAME}/{user_id}/pptx/tables/{csv_file_name}')

                        textdct["tables"] = tablelst
                    else:
                        textdct["tables"] = []

            except Exception as e:
                print("Unable to extract PPT tables, Error:", e)

            
            try: 
                for shape in slide.shapes:
                # Images - JPG

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: 
                        image = shape
                        image = image.image
                        image_bytes = image.blob

                        image_file_name = (f'{pptx_filename}_{shortuuid.uuid()}.jpg')

                        # download image
                        with open(f'{TEMP_DIR}/images/{image_file_name}', "wb") as imagefile:
                            imagefile.write(image_bytes)

                        # upload image to GCS bucket
                        upload_blob(bucket_name=BUCKET_NAME, 
                                source_file_name=f"{TEMP_DIR}/images/{image_file_name}", 
                                destination_blob_name=f'{user_id}/pptx/images/{image_file_name}')

                        imagelst.append(f'https://storage.googleapis.com/{BUCKET_NAME}/{user_id}/pptx/images/{image_file_name}')

                        textdct["images"] = imagelst
                    else:          
                        textdct["images"] = []

            except Exception as e:
                print("Unable to extract PPT images, Error:", e)

            slidename = "slide" + str(slidecounter)

            pptdct[slidename] = textdct

            slidecounter += 1

        # download json
        with open(f'{TEMP_DIR}/json/{json_file_name}', "w") as t:
            t.write('%s\n' % pptdct)

        # upload json to GCS bucket
        upload_blob(bucket_name=BUCKET_NAME, 
                source_file_name=f"{TEMP_DIR}/json/{json_file_name}", 
                destination_blob_name=f'{user_id}/pptx/json/{json_file_name}')    
        
        return pptdct
    except Exception as e:
                print("Unable to extract PPT file, Error:", e)
